"""Tests for the PowerPoint (.pptx) parser."""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from pith.parsers.base import ParseError
from pith.parsers.pptx import parse_pptx


def _make_pptx(path: Path, slides: list[dict[str, object]]) -> Path:
    """Create a .pptx file with given slide data.

    Each dict in *slides* may contain:
        title:  str — slide title text
        bodies: list[str] — text boxes to add
        table:  list[list[str]] — table data (first row = headers)
    """
    prs = Presentation()
    for spec in slides:
        layout = prs.slide_layouts[5]  # blank layout
        slide = prs.slides.add_slide(layout)

        if "title" in spec:
            # Add a title shape.
            from pptx.util import Pt
            txbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(1))
            txbox.text_frame.text = spec["title"]
            # Mark it as the title by using the title placeholder approach:
            # Instead, use a layout with title.
            prs.slides._sldIdLst.remove(prs.slides._sldIdLst[-1])
            slide_layout = prs.slide_layouts[1]  # title + content
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = spec["title"]

        for body_text in spec.get("bodies", []):
            txbox = slide.shapes.add_textbox(
                Inches(0.5), Inches(2), Inches(8), Inches(1),
            )
            txbox.text_frame.text = body_text

        table_data = spec.get("table")
        if table_data:
            rows = len(table_data)
            cols = len(table_data[0]) if table_data else 0
            tbl_shape = slide.shapes.add_table(
                rows, cols, Inches(0.5), Inches(4), Inches(8), Inches(2),
            )
            for r_idx, row in enumerate(table_data):
                for c_idx, val in enumerate(row):
                    tbl_shape.table.cell(r_idx, c_idx).text = val

    prs.save(str(path))
    return path


def _make_simple_pptx(path: Path) -> Path:
    """Create a simple PPTX with a titled slide and body text."""
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Introduction"
    body = slide.placeholders[1]
    body.text = "Welcome to the presentation"
    prs.save(str(path))
    return path


class TestParsePptx:
    def test_single_slide_with_title(self, tmp_path: Path) -> None:
        path = _make_simple_pptx(tmp_path / "simple.pptx")
        doc = parse_pptx(path)

        assert doc.title == "simple"
        assert len(doc.sections) == 1
        assert doc.sections[0].heading == "Introduction"
        assert "Welcome to the presentation" in doc.sections[0].body
        assert doc.metadata["source"] == str(path)
        assert doc.metadata["format"] == "pptx"
        assert doc.metadata["slide_count"] == "1"

    def test_slide_without_title_gets_fallback(self, tmp_path: Path) -> None:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank, no title
        txbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        txbox.text_frame.text = "Some content"
        path = tmp_path / "notitle.pptx"
        prs.save(str(path))

        doc = parse_pptx(path)

        assert len(doc.sections) == 1
        assert doc.sections[0].heading == "Slide 1"
        assert doc.sections[0].body == "Some content"

    def test_empty_slide_skipped(self, tmp_path: Path) -> None:
        prs = Presentation()
        # Slide 1: empty (blank layout, no text)
        prs.slides.add_slide(prs.slide_layouts[6])
        # Slide 2: has content
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "Real Slide"
        slide2.placeholders[1].text = "Content here"
        path = tmp_path / "with_empty.pptx"
        prs.save(str(path))

        doc = parse_pptx(path)

        assert len(doc.sections) == 1
        assert doc.sections[0].heading == "Real Slide"

    def test_table_extraction(self, tmp_path: Path) -> None:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Data Slide"
        tbl_shape = slide.shapes.add_table(
            3, 2, Inches(0.5), Inches(2), Inches(6), Inches(2),
        )
        tbl = tbl_shape.table
        tbl.cell(0, 0).text = "Name"
        tbl.cell(0, 1).text = "Score"
        tbl.cell(1, 0).text = "Alice"
        tbl.cell(1, 1).text = "95"
        tbl.cell(2, 0).text = "Bob"
        tbl.cell(2, 1).text = "87"
        path = tmp_path / "with_table.pptx"
        prs.save(str(path))

        doc = parse_pptx(path)

        assert len(doc.tables) == 1
        assert doc.tables[0].headers == ["Name", "Score"]
        assert doc.tables[0].rows == [["Alice", "95"], ["Bob", "87"]]

    def test_multiple_slides(self, tmp_path: Path) -> None:
        prs = Presentation()
        for i in range(3):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = f"Slide {i + 1} Title"
            slide.placeholders[1].text = f"Body {i + 1}"
        path = tmp_path / "multi.pptx"
        prs.save(str(path))

        doc = parse_pptx(path)

        assert len(doc.sections) == 3
        assert doc.sections[0].heading == "Slide 1 Title"
        assert doc.sections[2].heading == "Slide 3 Title"
        assert doc.metadata["slide_count"] == "3"

    def test_corrupt_file_raises(self, tmp_path: Path) -> None:
        path = tmp_path / "corrupt.pptx"
        path.write_bytes(b"not a pptx file")

        with pytest.raises(ParseError) as exc_info:
            parse_pptx(path)
        assert exc_info.value.path == path

    def test_metadata_shape_count(self, tmp_path: Path) -> None:
        path = _make_simple_pptx(tmp_path / "shapes.pptx")
        doc = parse_pptx(path)

        assert int(doc.metadata["shape_count"]) > 0

    def test_dispatch_integration(self, tmp_path: Path) -> None:
        """Verify .pptx routes through the dispatch map."""
        from pith.parsers import parse

        path = _make_simple_pptx(tmp_path / "dispatch.pptx")
        doc = parse(path)
        assert doc.metadata["format"] == "pptx"
