"""PPTX parser — extraction via python-pptx.

Each slide becomes a Section. Tables on slides are extracted
as Table objects.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from pith.i18n import t
from pith.parsers.base import ParsedDocument, ParseError, Section, Table


def _slide_title(slide: object, index: int) -> str:
    """Return the slide title text, or a fallback 'Slide N'."""
    if slide.shapes.title and slide.shapes.title.has_text_frame:
        text = slide.shapes.title.text_frame.text.strip()
        if text:
            return text
    return t("parser.pptx_slide_fallback", n=index)


def parse_pptx(path: Path) -> ParsedDocument:
    """Parse a PPTX file into a ``ParsedDocument``.

    Each slide becomes a section. Text from non-title shapes is
    joined into the section body. Tables are extracted separately.

    Args:
        path: Path to the PPTX file.

    Returns:
        A ``ParsedDocument`` with one section per slide.

    Raises:
        ParseError: If the file cannot be read.
    """
    try:
        prs = Presentation(str(path))
    except Exception as exc:
        raise ParseError(path, t("parser.pptx_read_error", detail=str(exc))) from exc

    sections: list[Section] = []
    tables: list[Table] = []
    total_shapes = 0

    for slide_index, slide in enumerate(prs.slides, start=1):
        total_shapes += len(slide.shapes)
        body_parts: list[str] = []
        slide_has_text = False

        for shape in slide.shapes: # type: ignore[attr-defined]
            # Skip the title shape — it becomes the heading.
            if shape is slide.shapes.title:
                if shape.has_text_frame:
                    title_text = shape.text_frame.text.strip()
                    if title_text:
                        slide_has_text = True
                continue

            if shape.has_table:
                tbl = shape.table
                rows_data: list[list[str]] = []
                for row in tbl.rows:
                    rows_data.append([cell.text.strip() for cell in row.cells])
                if rows_data:
                    tables.append(Table(headers=rows_data[0], rows=rows_data[1:]))
                    slide_has_text = True

            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        body_parts.append(text)
                        slide_has_text = True

        if not slide_has_text:
            continue

        heading = _slide_title(slide, slide_index)
        body = "\n".join(body_parts)
        sections.append(Section(heading=heading, body=body))

    return ParsedDocument(
        title=path.stem,
        sections=sections,
        tables=tables,
        metadata={
            "source": str(path),
            "format": "pptx",
            "slide_count": str(len(prs.slides)),
            "shape_count": str(total_shapes),
        },
    )
